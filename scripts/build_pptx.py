"""Apply a fill spec to a copy of master.pptx and write the output.

The spec is the authoritative instruction produced by the data-mapper
agent. This script is intentionally dumb: it copies the master template,
walks the spec, and replaces text in-place while preserving each run's
formatting (font, size, color, weight). It never adds, removes, or moves
shapes — that's how pixel fidelity to the template is preserved.

Spec format (JSON):
{
  "fills": [
    {"slide_idx": 0, "shape_idx": 5, "value": "..."},
    {"slide_idx": 1, "shape_idx": 2, "row": 3, "col": 4, "value": "..."}
  ]
}
"""
import argparse
import json
import shutil
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = ROOT / "templates" / "master.pptx"


def _replace_text_frame(tf, new_text: str) -> None:
    """Replace text content while preserving the first run's formatting.

    PPTX paragraphs hold runs; each run carries its own font/size/color.
    We keep the first run of the first paragraph, set its text to the new
    value, and clear every other run/paragraph. This loses run-level
    diversity inside the slot but preserves typography exactly.
    """
    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        tf.text = new_text
        return

    first = paragraphs[0]
    if first.runs:
        first.runs[0].text = new_text
        for r in first.runs[1:]:
            r.text = ""
    else:
        first.text = new_text

    for extra in paragraphs[1:]:
        extra._p.getparent().remove(extra._p)


def apply_spec(pres, spec: dict) -> list:
    issues = []
    for entry in spec.get("fills", []):
        s_idx = entry.get("slide_idx")
        sh_idx = entry.get("shape_idx")
        value = entry.get("value", "")
        try:
            shape = pres.slides[s_idx].shapes[sh_idx]
        except (IndexError, TypeError):
            issues.append(
                f"missing shape: slide={s_idx} shape={sh_idx}"
            )
            continue

        if "row" in entry and "col" in entry:
            if not shape.has_table:
                issues.append(
                    f"not a table: slide={s_idx} shape={sh_idx}"
                )
                continue
            try:
                cell = shape.table.cell(entry["row"], entry["col"])
            except Exception as e:
                issues.append(
                    f"bad cell ({entry['row']},{entry['col']}): {e}"
                )
                continue
            _replace_text_frame(cell.text_frame, str(value))
        else:
            if not shape.has_text_frame:
                issues.append(
                    f"no text frame: slide={s_idx} shape={sh_idx}"
                )
                continue
            _replace_text_frame(shape.text_frame, str(value))

    return issues


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--spec", required=True, help="path to fill spec JSON")
    ap.add_argument("--out", required=True, help="output pptx path")
    ap.add_argument("--template", default=str(DEFAULT_TEMPLATE))
    args = ap.parse_args()

    spec_path = Path(args.spec)
    out_path = Path(args.out)
    template_path = Path(args.template)

    if not template_path.exists():
        print(f"template not found: {template_path}", file=sys.stderr)
        sys.exit(2)
    if not spec_path.exists():
        print(f"spec not found: {spec_path}", file=sys.stderr)
        sys.exit(2)

    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(template_path, out_path)

    pres = Presentation(str(out_path))
    issues = apply_spec(pres, spec)
    pres.save(str(out_path))

    json.dump(
        {
            "output": str(out_path),
            "applied": len(spec.get("fills", [])) - len(issues),
            "issues": issues,
        },
        sys.stdout,
        ensure_ascii=False,
        indent=2,
    )


if __name__ == "__main__":
    main()
