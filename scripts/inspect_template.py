"""Dump master.pptx structure as JSON.

Used by the template-curator agent to learn what slots exist.
Output goes to stdout. Each slide lists every shape with a stable
(slide_idx, shape_idx) address; tables additionally list cells by (row, col).
The current text in each slot is included as a semantic hint for the
data-mapper agent.
"""
import json
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = ROOT / "templates" / "master.pptx"


def inspect(template_path: Path) -> dict:
    p = Presentation(str(template_path))
    out = {"template": str(template_path), "slides": []}
    for s_idx, slide in enumerate(p.slides):
        slide_info = {
            "slide_idx": s_idx,
            "layout": slide.slide_layout.name,
            "shapes": [],
        }
        for sh_idx, shp in enumerate(slide.shapes):
            entry = {
                "shape_idx": sh_idx,
                "name": shp.name,
                "type": str(shp.shape_type),
            }
            if shp.has_text_frame:
                entry["text"] = shp.text_frame.text
            if shp.has_table:
                cells = []
                for r_idx, row in enumerate(shp.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        cells.append(
                            {"row": r_idx, "col": c_idx, "text": cell.text}
                        )
                entry["table"] = {
                    "rows": len(shp.table.rows),
                    "cols": len(shp.table.columns),
                    "cells": cells,
                }
            slide_info["shapes"].append(entry)
        out["slides"].append(slide_info)
    return out


def main():
    template = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_TEMPLATE
    if not template.exists():
        print(f"template not found: {template}", file=sys.stderr)
        sys.exit(2)
    json.dump(inspect(template), sys.stdout, ensure_ascii=False, indent=2)


if __name__ == "__main__":
    main()
