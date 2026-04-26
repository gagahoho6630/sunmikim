"""Generate a generic 2-slide example template at templates/master.pptx.

Used to bootstrap the plugin without shipping any proprietary content.
Slide 1 = generic cover with title + overview table. Slide 2 = generic
5-day schedule table. The agents work on whatever shape structure exists
in master.pptx, so users can replace this file with their own template
and the pipeline continues to function.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "templates" / "master.pptx"

NAVY = RGBColor(0x12, 0x2E, 0x5A)
LIGHT = RGBColor(0xE8, 0xEE, 0xF7)
DARK = RGBColor(0x22, 0x22, 0x22)


def add_textbox(slide, left, top, width, height, text, *, size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Pretendard"
    return box


def style_table(table, header_fill=NAVY, header_color=RGBColor(0xFF, 0xFF, 0xFF)):
    for c_idx, cell in enumerate(table.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = header_color
                run.font.bold = True
                run.font.size = Pt(10)
                run.font.name = "Pretendard"
    for r_idx in range(1, len(table.rows)):
        for cell in table.rows[r_idx].cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)
                    run.font.name = "Pretendard"


def build_cover(slide):
    add_textbox(
        slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.7),
        "Project Title — Replace via fill spec",
        size=22, bold=True, color=NAVY,
    )
    add_textbox(
        slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5),
        "1. Overview", size=14, bold=True, color=NAVY,
    )
    add_textbox(
        slide, Inches(0.5), Inches(1.7), Inches(9), Inches(1.2),
        "가. 사업명: ...\n나. 계약기간: ...\n다. 운영장소: ...\n라. 대상: ...",
        size=11,
    )
    rows, cols = 3, 3
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(3.1), Inches(9), Inches(2.0),
    ).table
    headers = ["Category", "Detail", "Note"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    sample = [
        ["Period", "YYYY.MM ~ YYYY.MM", ""],
        ["Audience", "...", ""],
    ]
    for r, row in enumerate(sample, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    style_table(table)

    add_textbox(
        slide, Inches(0.5), Inches(5.3), Inches(9), Inches(0.5),
        "2. Submission", size=14, bold=True, color=NAVY,
    )
    add_textbox(
        slide, Inches(0.5), Inches(5.8), Inches(9), Inches(1.5),
        "Replace this paragraph with submission details, schedule, and contact.",
        size=11,
    )


def build_schedule(slide):
    add_textbox(
        slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.7),
        "[Attachment] Sample 5-Day Schedule",
        size=20, bold=True, color=NAVY,
    )
    rows, cols = 11, 6
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.3), Inches(9), Inches(5.5),
    ).table
    headers = ["Time", "Mon (Day 1)", "Tue (Day 2)", "Wed (Day 3)", "Thu (Day 4)", "Fri (Day 5)"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    times = [
        "08:00-09:00", "09:00-10:00", "10:00-11:00", "11:00-12:00",
        "12:00-13:00", "13:00-14:00", "14:00-15:00", "15:00-16:00",
        "16:00-17:00", "17:00-18:00",
    ]
    for r, t in enumerate(times, 1):
        table.cell(r, 0).text = t
    style_table(table)


def main():
    prs = Presentation()
    prs.slide_width = Emu(9906000)
    prs.slide_height = Emu(6858000)

    blank = prs.slide_layouts[6]
    build_cover(prs.slides.add_slide(blank))
    build_schedule(prs.slides.add_slide(blank))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
